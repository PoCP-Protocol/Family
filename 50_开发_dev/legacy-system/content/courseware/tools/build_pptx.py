"""deck-generation Skill(.pptx)+ instructional-design 母版。
课件 YAML → 真 PowerPoint(python-pptx)。品牌母版:配色与 HTML deck 统一(橙/蓝),阶段色带,低认知负荷 4 卡布局。
用法: python build_pptx.py [courseware.yaml]   默认 ../parent_21day_camp.yaml
输出: ../decks/<code>/<code>.pptx
"""
from __future__ import annotations
import sys
from pathlib import Path
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
SRC = Path(sys.argv[1]) if len(sys.argv) > 1 else HERE.parent / "parent_21day_camp.yaml"

# ---- family 品牌母版(与 build_deck.py 调色板一致)----
ORANGE = RGBColor(0xEA, 0x73, 0x17)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
INK    = RGBColor(0x1F, 0x29, 0x37)
GREY   = RGBColor(0x6B, 0x72, 0x80)
BG_A   = RGBColor(0xFF, 0xF7, 0xED)  # 暖(今日一件小事)
BG_B   = RGBColor(0xEF, 0xF6, 0xFF)  # 蓝(今晚可说的话)
BG_C   = RGBColor(0xF8, 0xFA, 0xFC)  # 观察
BG_D   = RGBColor(0xF9, 0xFA, 0xFB)  # 边界
W, H = Inches(13.333), Inches(7.5)
# 阶段配色带(与 HTML deck / 设计系统一致)
TEAL = RGBColor(0x0F, 0x76, 0x6E); VIOLET = RGBColor(0x7C, 0x3A, 0xED)
PHASE_ACCENT = {"SEE_CONNECT": ORANGE, "SEE": ORANGE, "PARENT_FIRST": BLUE,
                "CO_CREATE": TEAL, "CO_STABILIZE": TEAL, "STABILIZE": VIOLET}


def is_review(d) -> bool:
    s = str(d.get("skill", "")) + str(d.get("theme", ""))
    return "复盘" in s or "回看" in s


def box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill or RGBColor(0xFF, 0xFF, 0xFF)
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (s, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Microsoft YaHei"
    return tb


def card(slide, l, t, w, h, title, body, bg, accent):
    box(slide, l, t, w, h, fill=bg, line=accent)
    tb = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12), w - Inches(0.36), h - Inches(0.24))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(12); r.font.color.rgb = GREY; r.font.bold = True; r.font.name = "Microsoft YaHei"
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(14); r2.font.color.rgb = INK; r2.font.name = "Microsoft YaHei"


def build(doc: dict, prs: Presentation):
    prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]
    m = doc.get("meta", {})
    ev = {e["id"]: e for e in m.get("evidence_grounding", [])}
    phases = {p["code"]: p for p in m.get("phases", [])}

    def ev_text(refs):
        if not refs: return "依据:practitioner · 证据待验"
        out = []
        for r in refs:
            e = ev.get(r)
            out.append(f'{e.get("grade","")} doi:{e["doi"]}' if e else str(r))
        return "依据:" + " ; ".join(out)

    # 封面
    s = prs.slides.add_slide(blank)
    box(s, 0, 0, W, Inches(0.35), fill=ORANGE)
    txt(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.5), [(m.get("camp_code",""), 16, ORANGE, True)])
    txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.4), [(m.get("title",""), 40, INK, True)])
    txt(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6), [(m.get("philosophy",""), 20, INK, False)])
    ph = " · ".join(f'{p["title"]}(D{p["day_from"]}–{p["day_to"]})' for p in m.get("phases", []))
    txt(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6), [(ph, 15, GREY, False)])
    txt(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.5), [("对象:" + m.get("audience",""), 14, GREY, False)])

    lessons = doc.get("days") or doc.get("weeks") or []
    unit = "Day" if doc.get("days") else "Week"
    total = len(lessons)

    def foot(s, num):
        txt(s, Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.3), [(f'{m.get("camp_code","")}    {unit} {num} / {total}', 9, RGBColor(0xB8,0xC0,0xCC), False)])

    seen_phase = set()
    for d in lessons:
        pcode = d.get("phase", ""); accent = PHASE_ACCENT.get(pcode, ORANGE)
        num = d.get("day", d.get("week"))
        ptitle = phases.get(pcode, {}).get("title", "")
        # 阶段分隔页
        if pcode and pcode not in seen_phase:
            seen_phase.add(pcode); p = phases.get(pcode, {})
            s = prs.slides.add_slide(blank)
            box(s, 0, 0, W, Inches(0.35), fill=accent)
            txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5), [(f"阶段 {len(seen_phase)} / {len(phases)}", 15, GREY, True)])
            txt(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2), [(ptitle, 46, accent, True)])
            txt(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5), [(f'D{p.get("day_from","")}–{p.get("day_to","")}', 18, GREY, False)])
        # 复盘/结营日:单独版式(一页一观点,无 2x2)
        if is_review(d):
            is_last = num == total
            s = prs.slides.add_slide(blank)
            box(s, 0, 0, W, Inches(0.25), fill=accent)
            txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.5), [(f'{unit} {num}   ·   {ptitle} · {"结营" if is_last else "复盘"}', 14, GREY, True)])
            txt(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.8), [(d.get("theme",""), 30, INK, True)])
            txt(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.7), [(d.get("why",""), 14, GREY, False)])
            card(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(2.6), "回看一件事", str(d.get("parent_action","")) + "\n\n可以说:“" + str(d.get("say_it_tonight","")) + "”", BG_C, accent)
            txt(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.5), [(str(d.get("boundary","")), 12, GREY, False)])
            foot(s, num); continue
        # 日常课:4 卡 2x2
        s = prs.slides.add_slide(blank)
        box(s, 0, 0, W, Inches(0.25), fill=accent)
        txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.5), [(f"{unit} {num}   ·   {ptitle}", 14, GREY, True)])
        txt(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.8), [(d.get("theme",""), 30, INK, True)])
        txt(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.4), [("技能:" + str(d.get("skill","")), 15, accent, True)])
        txt(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.7), [(d.get("why",""), 13, GREY, False)])
        micro = d.get("parent_daily_micro")
        a_title = "本周焦点 + 每日微行动" if micro else "今日一件小事"
        a_body = (str(d.get("weekly_focus","")) + "\n" + "\n".join("· " + str(x) for x in micro)) if micro else str(d.get("parent_action",""))
        cw, ch = Inches(5.7), Inches(1.7); x1, x2 = Inches(0.9), Inches(6.9); y1, y2 = Inches(3.0), Inches(4.85)
        card(s, x1, y1, cw, ch, a_title, a_body, BG_A, ORANGE)
        card(s, x2, y1, cw, ch, "今晚可以说的话", "“" + str(d.get("say_it_tonight","")) + "”", BG_B, BLUE)
        card(s, x1, y2, cw, ch, "观察什么(过程,非评分)", str(d.get("look_for","")), BG_C, RGBColor(0xE5,0xE7,0xEB))
        card(s, x2, y2, cw, ch, "边界 / 提醒", str(d.get("boundary","")), BG_D, RGBColor(0xE5,0xE7,0xEB))
        txt(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.35), [(ev_text(d.get("evidence_refs")), 10, GREY, False)])
        foot(s, num)

    # 依据页
    s = prs.slides.add_slide(blank)
    box(s, 0, 0, W, Inches(0.25), fill=BLUE)
    txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.6), [("循证依据(crossref 已核验)", 26, INK, True)])
    runs = [(f'{e.get("grade","")} · doi:{e["doi"]} — {e.get("note","")}', 14, INK, False) for e in m.get("evidence_grounding", [])]
    txt(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(3.5), runs or [("(无)", 14, GREY, False)])
    txt(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8), [("红线:" + " · ".join(m.get("red_lines", [])), 12, ORANGE, False)])


def main() -> int:
    doc = yaml.safe_load(SRC.read_text(encoding="utf-8"))
    code = doc.get("meta", {}).get("camp_code", "COURSE")
    out_dir = HERE.parent / "decks" / code; out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(); build(doc, prs)
    out = out_dir / f"{code}.pptx"; prs.save(str(out))
    print(f"pptx built: {out}  slides={len(prs.slides._sldIdLst)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
